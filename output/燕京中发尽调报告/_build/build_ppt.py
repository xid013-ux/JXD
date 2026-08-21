# -*- coding: utf-8 -*-
"""燕京中发业务架构调整方案——调整前后对比（单页PPT）"""
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE

BLUE   = RGBColor(0x1F, 0x9B, 0xDE)   # 主体框（参考图配色）
BLUE_D = RGBColor(0x15, 0x6E, 0x9E)   # 深蓝：上市主体
GREY   = RGBColor(0xB7, 0xD6, 0xE4)   # 浅色：非核心/待处置
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x33, 0x33, 0x33)
LINE   = RGBColor(0x59, 0x59, 0x59)
RED    = RGBColor(0xC0, 0x3A, 0x2B)

prs = Presentation()
prs.slide_width, prs.slide_height = Cm(33.87), Cm(19.05)   # 16:9
s = prs.slides.add_slide(prs.slide_layouts[6])

def box(x, y, w, h, text, fill=BLUE, fc=WHITE, sz=10.5, bold=True, line=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(1)
    sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Cm(0.1)
    tf.margin_top = tf.margin_bottom = Cm(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    for i, ln in enumerate(text.split('\n')):
        pp = p if i == 0 else tf.add_paragraph()
        pp.alignment = PP_ALIGN.CENTER
        r = pp.add_run(); r.text = ln
        r.font.size = Pt(sz); r.font.bold = bold
        r.font.color.rgb = fc; r.font.name = '微软雅黑'
    return sh

def label(x, y, w, h, text, sz=9, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Cm(0.05)
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    for i, ln in enumerate(text.split('\n')):
        pp = p if i == 0 else tf.add_paragraph()
        pp.alignment = align
        r = pp.add_run(); r.text = ln
        r.font.size = Pt(sz); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = '微软雅黑'
    return tb

def vline(x, y1, y2, arrow=True, color=LINE):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(x), Cm(y1), Cm(x), Cm(y2))
    c.line.color.rgb = color; c.line.width = Pt(1.25)
    if arrow:
        c.line._get_or_add_ln().append(__import__('lxml.etree', fromlist=['etree']).fromstring(
            '<a:tailEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" type="triangle" w="med" len="med"/>'))
    return c

def hline(x1, x2, y, color=LINE):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(x1), Cm(y), Cm(x2), Cm(y))
    c.line.color.rgb = color; c.line.width = Pt(1.25)
    return c

# ================= 标题 =================
label(1.0, 0.55, 32, 1.0, '燕京中发业务架构调整方案——架构调整前后对比', sz=18, bold=True, color=RGBColor(0x15,0x3E,0x5C))
ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(1.0), Cm(1.62), Cm(32.6), Cm(1.62))
ln.line.color.rgb = BLUE; ln.line.width = Pt(2)

label(1.0, 1.75, 14.5, 0.7, '一、调整前', sz=13, bold=True, color=RGBColor(0x15,0x3E,0x5C))
label(19.4, 1.75, 13.5, 0.7, '二、调整后', sz=13, bold=True, color=RGBColor(0x15,0x3E,0x5C))

# ================= 左：调整前 =================
box(1.0, 2.7, 4.0, 1.05, '燕京啤酒\n（000729）', sz=10)
box(5.6, 2.7, 3.5, 1.05, '燕京集团', sz=10)
box(9.7, 2.7, 4.2, 1.05, '食品发酵\n研究院', sz=10)

# 三股东 -> 发行人
hline(3.0, 11.8, 4.35)
vline(3.0, 3.75, 4.35, arrow=False); vline(7.35, 3.75, 4.35, arrow=False); vline(11.8, 3.75, 4.35, arrow=False)
label(3.15, 3.80, 2.0, 0.5, '80%', sz=9, bold=True)
label(7.50, 3.80, 2.2, 0.5, '17.5%', sz=9, bold=True)
label(11.95, 3.80, 2.0, 0.5, '2.5%', sz=9, bold=True)
vline(7.35, 4.35, 5.05)
box(3.6, 5.05, 7.5, 1.15, '北京燕京中发生物技术有限公司（发行人）', fill=BLUE_D, sz=10.5)

# 发行人 -> 中发邢台 33%
vline(7.35, 6.20, 7.55)
label(7.50, 6.55, 2.0, 0.5, '33%', sz=9, bold=True)
box(4.3, 7.55, 6.1, 1.15, '燕京中发生物技术\n（邢台）有限公司', sz=10)

# 燕京啤酒 -> 中发邢台 67%
vline(1.6, 3.75, 8.13, arrow=False)
hline(1.6, 4.3, 8.13)
c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(1.6), Cm(8.13), Cm(4.3), Cm(8.13))
c.line.color.rgb = LINE; c.line.width = Pt(1.25)
label(1.75, 7.62, 2.0, 0.5, '67%', sz=9, bold=True)

# 燕京啤酒 -> 燕京啤酒（邢台）100%
vline(2.6, 3.75, 10.6, arrow=False)
hline(2.6, 4.3, 10.6)
c2 = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(2.6), Cm(10.6), Cm(4.3), Cm(10.6))
c2.line.color.rgb = LINE; c2.line.width = Pt(1.25)
label(2.75, 10.08, 2.2, 0.5, '100%', sz=9, bold=True)
box(4.3, 10.05, 6.1, 1.15, '燕京啤酒（邢台）\n有限公司', fill=GREY, fc=RGBColor(0x15,0x3E,0x5C), sz=10)
label(4.3, 11.30, 8.6, 2.2,
      '· 持有四宗宗地，合计 181,900.00 ㎡（272.85 亩）\n'
      '· 出让 3 宗 215.50 亩；集体建设用地 1 宗 57.35 亩\n'
      '· 厂房已整体转为纳豆生产，不再从事啤酒业务',
      sz=8.5, color=RGBColor(0x55,0x55,0x55))

label(1.0, 13.9, 13.3, 2.6,
      '现状问题：\n'
      '① 中发邢台与发行人主营同为纳豆，构成同业竞争；\n'
      '② 生产用地及厂房登记在关联方名下，房、地权属分离，资产完整性存疑。',
      sz=9, color=RED)

# ================= 中间：调整步骤 =================
arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Cm(14.9), Cm(6.6), Cm(3.9), Cm(1.5))
arrow.fill.solid(); arrow.fill.fore_color.rgb = BLUE
arrow.line.fill.background(); arrow.shadow.inherit = False
arrow.text_frame.text = ''

label(14.45, 8.4, 4.9, 5.4,
      '调整步骤\n'
      '1、发行人以自身股权为对价，收购燕京啤酒持有的中发邢台 67% 股权\n'
      '2、发行人以自身股权为对价，收购燕京啤酒持有的燕京啤酒（邢台）100% 股权\n'
      '3、标的公司更名、变更经营范围，注销啤酒生产许可\n'
      '4、两家邢台主体择机吸收合并，实现房、地权属统一',
      sz=8.5, color=RGBColor(0x15,0x3E,0x5C))

# ================= 右：调整后 =================
box(19.4, 2.7, 4.0, 1.05, '燕京啤酒\n（000729）', sz=10)
box(24.0, 2.7, 3.5, 1.05, '燕京集团', sz=10)
box(28.1, 2.7, 4.2, 1.05, '食品发酵\n研究院', sz=10)

hline(21.4, 30.2, 4.35)
vline(21.4, 3.75, 4.35, arrow=False); vline(25.75, 3.75, 4.35, arrow=False); vline(30.2, 3.75, 4.35, arrow=False)
label(21.55, 3.80, 2.4, 0.5, '↑ 上升', sz=8.5, bold=True, color=RED)
label(25.90, 3.80, 2.4, 0.5, '↓ 稀释', sz=8.5, bold=True, color=RGBColor(0x88,0x88,0x88))
label(30.35, 3.80, 2.4, 0.5, '↓ 稀释', sz=8.5, bold=True, color=RGBColor(0x88,0x88,0x88))
vline(25.75, 4.35, 5.05)
box(22.0, 5.05, 7.5, 1.15, '北京燕京中发生物技术有限公司（发行人）', fill=BLUE_D, sz=10.5)

# 发行人 -> 两家子公司
hline(22.4, 29.6, 7.05)
vline(25.75, 6.20, 7.05, arrow=False)
vline(22.4, 7.05, 7.75); vline(29.6, 7.05, 7.75)
label(22.55, 7.08, 2.2, 0.5, '100%', sz=9, bold=True)
label(27.75, 7.08, 2.2, 0.5, '100%', sz=9, bold=True)
box(19.5, 7.75, 5.8, 1.35, '燕京中发生物技术\n（邢台）有限公司', sz=10)
box(26.6, 7.75, 6.0, 1.35, '燕京啤酒（邢台）有限公司\n（更名，去除“燕京啤酒”字样）', sz=9.5)

# 吸收合并
merge = s.shapes.add_shape(MSO_SHAPE.LEFT_RIGHT_ARROW, Cm(25.45), Cm(8.15), Cm(1.05), Cm(0.6))
merge.fill.solid(); merge.fill.fore_color.rgb = RGBColor(0xE8,0xA0,0x33)
merge.line.fill.background(); merge.shadow.inherit = False
label(23.5, 9.25, 4.5, 0.5, '择机吸收合并', sz=8.5, bold=True, color=RGBColor(0xB0,0x76,0x18), align=PP_ALIGN.CENTER)

vline(25.75, 9.85, 10.5)
box(21.0, 10.5, 9.5, 1.15, '合并后主体（持有土地、厂房及纳豆业务）', fill=BLUE_D, sz=10.5)
label(21.0, 11.75, 11.6, 1.8,
      '· 四宗宗地及地上厂房、在建工程全部纳入上市主体\n'
      '· 房、地权属统一，资产完整性、独立性问题一并解决',
      sz=8.5, color=RGBColor(0x55,0x55,0x55))

label(19.4, 13.9, 13.5, 2.6,
      '调整效果：\n'
      '① 中发邢台成为全资子公司并表，同业竞争消除；\n'
      '② 土地、厂房随股权进入上市主体，不涉及不动产过户，不产生土地增值税及契税；\n'
      '③ 股权支付比例 100%，可争取企业所得税特殊性税务处理，预计仅产生印花税。',
      sz=9, color=RGBColor(0x1B,0x7A,0x3D))

# ================= 页脚 =================
foot = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(1.0), Cm(17.5), Cm(32.6), Cm(17.5))
foot.line.color.rgb = RGBColor(0xCC,0xCC,0xCC); foot.line.width = Pt(0.75)
label(1.0, 17.65, 24, 0.9,
      '注：调整后持股比例以经国资监管部门备案的评估结果为准；本图为方案示意。',
      sz=8, color=RGBColor(0x88,0x88,0x88))
label(26.0, 17.65, 6.6, 0.9, '国泰海通证券  燕京中发IPO项目组', sz=8,
      color=RGBColor(0x88,0x88,0x88), align=PP_ALIGN.RIGHT)

import os
out = '/home/user/JXD/output/燕京中发尽调报告/燕京中发业务架构调整方案-架构调整前后对比.pptx'
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print('已生成:', out)
print('页面 %.2f x %.2f cm，形状数 %d' % (prs.slide_width/360000, prs.slide_height/360000, len(s.shapes)))
